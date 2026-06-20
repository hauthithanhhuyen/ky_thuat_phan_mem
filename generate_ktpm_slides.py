import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# --- Color Palette ---
COLOR_BG_DARK = RGBColor(15, 23, 42)        # Slate 900 (Dark Slate)
COLOR_BG_LIGHT = RGBColor(248, 250, 252)    # Slate 50 (Off-white)
COLOR_PRIMARY_TEXT = RGBColor(15, 23, 42)   # Slate 900 (Title text)
COLOR_SECONDARY_TEXT = RGBColor(71, 85, 105) # Slate 600 (Body text)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ACCENT_BLUE = RGBColor(37, 99, 235)    # Blue 600 (Primary accent)
COLOR_ACCENT_TEAL = RGBColor(13, 148, 136)   # Teal 600 (Secondary accent)
COLOR_BORDER_LIGHT = RGBColor(226, 232, 240) # Slate 200 (Line separators)

# Soft Tint Colors for Cards
CARD_INFO_BG = RGBColor(240, 246, 255)       # Light Blue
CARD_INFO_BORDER = RGBColor(59, 130, 246)
CARD_SUCCESS_BG = RGBColor(240, 253, 250)    # Light Teal
CARD_SUCCESS_BORDER = RGBColor(13, 148, 136)
CARD_WARN_BG = RGBColor(254, 242, 242)       # Light Rose
CARD_WARN_BORDER = RGBColor(225, 29, 72)
CARD_DEFAULT_BG = RGBColor(255, 255, 255)    # Plain White
CARD_DEFAULT_BORDER = RGBColor(226, 232, 240)

# --- Drawing Helpers ---
def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def draw_accent_bar(slide, dark_theme=False):
    # Left edge vertical accent line
    accent_color = COLOR_ACCENT_TEAL if dark_theme else COLOR_ACCENT_BLUE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color
    line.line.fill.solid()
    line.line.fill.fore_color.rgb = accent_color

def draw_horizontal_line(slide, left, top, width, height, color):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.solid()
    line.line.fill.fore_color.rgb = color
    return line

def add_footer(slide, slide_num):
    # Separator line
    draw_horizontal_line(slide, Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.015), COLOR_BORDER_LIGHT)
    
    # Left Footer text
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(8.0), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Tiểu luận Kỹ thuật Phần mềm | Hầu Thị Thanh Huyền - K58.KTP"
    p.font.name = "Segoe UI"
    p.font.size = Pt(9.5)
    p.font.color.rgb = COLOR_SECONDARY_TEXT
    
    # Right Slide Number
    txBox2 = slide.shapes.add_textbox(Inches(11.5), Inches(6.9), Inches(1.0), Inches(0.4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = f"Trang {slide_num}"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = COLOR_SECONDARY_TEXT
    p2.alignment = PP_ALIGN.RIGHT

def parse_and_add_text(tf, text, font_size=16, space_after=14, dark_theme=False):
    """
    Format helper that parses '**Heading:** Description' and renders them
    as separate paragraphs.
    - No bullet points (dots).
    - Font sizes are increased.
    - Header is 4pt larger than Description.
    """
    primary_color = COLOR_WHITE if dark_theme else COLOR_PRIMARY_TEXT
    secondary_color = COLOR_WHITE if dark_theme else COLOR_SECONDARY_TEXT
    
    if text.startswith("**") and "**" in text[2:]:
        end_idx = text.find("**", 2)
        bold_part = text[2:end_idx]
        normal_part = text[end_idx+2:]
        
        # Paragraph for the Bold Heading (Larger Text)
        p_bold = tf.add_paragraph() if len(tf.paragraphs[0].text) > 0 else tf.paragraphs[0]
        r1 = p_bold.add_run()
        r1.text = bold_part
        r1.font.name = "Segoe UI"
        r1.font.size = Pt(font_size + 4)  # Bold headers are made larger (e.g. 20pt for 16pt body)
        r1.font.bold = True
        r1.font.color.rgb = primary_color
        p_bold.space_after = Pt(4)        # Tight space before its description
        
        # Paragraph for the normal description text
        p_norm = tf.add_paragraph()
        r2 = p_norm.add_run()
        r2.text = normal_part.strip()
        r2.font.name = "Segoe UI"
        r2.font.size = Pt(font_size)      # Description is base font size (e.g. 16pt)
        r2.font.bold = False
        r2.font.color.rgb = secondary_color
        p_norm.space_after = Pt(space_after) # Space before the next block
    else:
        # Regular text block
        p = tf.add_paragraph() if len(tf.paragraphs[0].text) > 0 else tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = "Segoe UI"
        r.font.size = Pt(font_size)
        r.font.bold = False
        r.font.color.rgb = primary_color
        p.space_after = Pt(space_after)

# --- Template Functions ---
def create_title_slide(title, subtitle, info_items):
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_BG_DARK)
    draw_accent_bar(slide, dark_theme=True)
    
    # Large Decorative Card/Accent background block
    accent_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.08))
    accent_box.fill.solid()
    accent_box.fill.fore_color.rgb = COLOR_ACCENT_TEAL
    accent_box.line.fill.solid()
    accent_box.line.fill.fore_color.rgb = COLOR_ACCENT_TEAL
    
    # Title Text Box - Font: Georgia (premium editorial/academic look)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Georgia"
    p.font.size = Pt(38)  # Increased from 36
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.name = "Georgia"
    p2.font.size = Pt(20)  # Increased from 18
    p2.font.color.rgb = COLOR_ACCENT_TEAL
    
    # Presenter Information box (bottom)
    info_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(11.733), Inches(2.2))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    
    for item in info_items:
        p_info = tf_info.add_paragraph() if len(tf_info.paragraphs[0].text) > 0 else tf_info.paragraphs[0]
        p_info.space_after = Pt(6)
        
        if ":" in item:
            label, val = item.split(":", 1)
            r1 = p_info.add_run()
            r1.text = label + ":"
            r1.font.name = "Segoe UI"
            r1.font.size = Pt(14)  # Increased from 13.5
            r1.font.bold = True
            r1.font.color.rgb = COLOR_ACCENT_TEAL
            
            r2 = p_info.add_run()
            r2.text = val
            r2.font.name = "Segoe UI"
            r2.font.size = Pt(14)
            r2.font.color.rgb = COLOR_WHITE
        else:
            r = p_info.add_run()
            r.text = item
            r.font.name = "Segoe UI"
            r.font.size = Pt(14)
            r.font.color.rgb = COLOR_WHITE
            
    return slide

def create_divider_slide(chapter_num, chapter_title):
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_BG_DARK)
    draw_accent_bar(slide, dark_theme=True)
    
    # Chapter tag
    txBox1 = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.0), Inches(0.6))
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = chapter_num.upper()
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(22)  # Increased from 20
    p1.font.bold = True
    p1.font.color.rgb = COLOR_ACCENT_TEAL
    
    # Accent Line
    draw_horizontal_line(slide, Inches(1.5), Inches(2.8), Inches(2.5), Inches(0.06), COLOR_ACCENT_TEAL)
    
    # Chapter Title - Font: Georgia (premium editorial/academic look)
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10.0), Inches(2.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = chapter_title
    p2.font.name = "Georgia"
    p2.font.size = Pt(40)  # Increased from 38
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.space_before = Pt(10)
    
    return slide

def create_content_slide(slide_num, section_tag, title, bullets=None, layout_type="single", columns_data=None, image_path=None, image_rect=None):
    """
    layout_type: "single", "split", "split_image"
    """
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_BG_LIGHT)
    draw_accent_bar(slide, dark_theme=False)
    add_footer(slide, slide_num)
    
    # Top Tag
    txBoxTag = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.4))
    tfTag = txBoxTag.text_frame
    pTag = tfTag.paragraphs[0]
    pTag.text = section_tag.upper()
    pTag.font.name = "Segoe UI"
    pTag.font.size = Pt(11)  # Increased from 10
    pTag.font.bold = True
    pTag.font.color.rgb = COLOR_ACCENT_BLUE
    
    # Slide Title - Font: Georgia (elegant, high-contrast style)
    txBoxTitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.733), Inches(0.8))
    tfTitle = txBoxTitle.text_frame
    pTitle = tfTitle.paragraphs[0]
    pTitle.text = title
    pTitle.font.name = "Georgia"
    pTitle.font.size = Pt(30)  # Increased from 26 (Large, bold title)
    pTitle.font.bold = True
    pTitle.font.color.rgb = COLOR_PRIMARY_TEXT
    
    # Accent line below title
    draw_horizontal_line(slide, Inches(0.8), Inches(1.55), Inches(11.733), Inches(0.025), COLOR_ACCENT_TEAL)
    
    # Content Area - Start from Top=1.8
    if layout_type == "single" and bullets:
        txBoxContent = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.8))
        tfContent = txBoxContent.text_frame
        tfContent.word_wrap = True
        tfContent.margin_left = Inches(0)
        tfContent.margin_top = Inches(0)
        
        # Single column layout has larger base text size: 16pt body, 20pt bold titles
        for bullet in bullets:
            parse_and_add_text(tfContent, bullet, font_size=16, space_after=16)
            
    elif layout_type == "split" and columns_data:
        # Col 1
        txBoxCol1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
        tfCol1 = txBoxCol1.text_frame
        tfCol1.word_wrap = True
        tfCol1.margin_left = Inches(0)
        tfCol1.margin_top = Inches(0)
        for bullet in columns_data[0]:
            parse_and_add_text(tfCol1, bullet, font_size=15, space_after=14)
            
        # Col 2
        txBoxCol2 = slide.shapes.add_textbox(Inches(6.933), Inches(1.8), Inches(5.6), Inches(4.8))
        tfCol2 = txBoxCol2.text_frame
        tfCol2.word_wrap = True
        tfCol2.margin_left = Inches(0)
        tfCol2.margin_top = Inches(0)
        for bullet in columns_data[1]:
            parse_and_add_text(tfCol2, bullet, font_size=15, space_after=14)
            
    elif layout_type == "split_image" and bullets and image_path:
        # Left Text Column
        txBoxCol1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
        tfCol1 = txBoxCol1.text_frame
        tfCol1.word_wrap = True
        tfCol1.margin_left = Inches(0)
        tfCol1.margin_top = Inches(0)
        for bullet in bullets:
            parse_and_add_text(tfCol1, bullet, font_size=14.5, space_after=14)
            
        # Right Image Position
        img_left, img_top, img_w, img_h = image_rect if image_rect else (Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5))
        
        # Background card card behind image
        bg_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_left - Inches(0.1), img_top - Inches(0.1), img_w + Inches(0.2), img_h + Inches(0.2))
        bg_card.fill.solid()
        bg_card.fill.fore_color.rgb = COLOR_WHITE
        bg_card.line.fill.solid()
        bg_card.line.fill.fore_color.rgb = COLOR_BORDER_LIGHT
        bg_card.line.width = Pt(1)
        
        slide.shapes.add_picture(image_path, img_left, img_top, img_w, img_h)
            
    return slide

def draw_card(slide, left, top, width, height, card_title, bullets, card_type="default"):
    if card_type == "info":
        bg_col, border_col = CARD_INFO_BG, CARD_INFO_BORDER
    elif card_type == "success":
        bg_col, border_col = CARD_SUCCESS_BG, CARD_SUCCESS_BORDER
    elif card_type == "warn":
        bg_col, border_col = CARD_WARN_BG, CARD_WARN_BORDER
    else:
        bg_col, border_col = CARD_DEFAULT_BG, CARD_DEFAULT_BORDER
        
    # Draw Background Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_col
    box.line.fill.solid()
    box.line.fill.fore_color.rgb = border_col
    box.line.width = Pt(1.5)
    
    # Draw Left Edge Accent Line
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.08), Inches(0.08), height - Inches(0.16))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = border_col
    accent_bar.line.fill.solid()
    accent_bar.line.fill.fore_color.rgb = border_col
    
    # Text Box inside card
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.12), width - Inches(0.35), height - Inches(0.24))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    # Title
    p_title = tf.paragraphs[0]
    p_title.text = card_title
    p_title.font.name = "Segoe UI"
    p_title.font.size = Pt(14.5)  # Increased from 14
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_PRIMARY_TEXT
    p_title.space_after = Pt(8)
    
    # Bullets inside card (Card content size set to 13pt)
    for bullet in bullets:
        parse_and_add_text(tf, bullet, font_size=13, space_after=6)

# --- Generate Slides ---

# Slide 1: Title
create_title_slide(
    "NỀN TẢNG AI/ML CẦN BIẾT CHO KIỂM THỬ PHẦN MỀM HIỆN ĐẠI",
    "Ứng dụng AI/ML trong kiểm thử và thực nghiệm phân loại mức độ nghiêm trọng của lỗi",
    [
        "Giáo viên hướng dẫn: TS. NGUYỄN TUẤN LINH",
        "Sinh viên thực hiện: HẦU THỊ THANH HUYỀN",
        "Mã số sinh viên: K225480106027",
        "Lớp: K58.KTP (Khoa Điện Tử)",
        "Trường: Đại học Kỹ thuật Công nghiệp Thái Nguyên (2026)"
    ]
)

# Slide 2: Lời nói đầu
create_content_slide(
    2,
    "Giới thiệu",
    "Lời nói đầu & Bối cảnh đề tài",
    bullets=[
        "**Chuyển đổi số mạnh mẽ:** Phần mềm ngày càng đóng vai trò cốt lõi trong đời sống và sản xuất. Đảm bảo chất lượng sản phẩm là yếu tố sống còn của một hệ thống.",
        "**Thách thức của QA truyền thống:** Quy trình kiểm thử thủ công và viết kịch bản lặp đi lặp lại tốn nhiều thời gian, công sức và dễ sai sót khi hệ thống mở rộng.",
        "**Xu hướng công nghệ mới:** Sự bùng nổ của Trí tuệ nhân tạo (AI) và Học máy (Machine Learning - ML) mở ra các hướng tiếp cận tự động hóa thông minh.",
        "**Mục đích đề tài:** Nghiên cứu lý thuyết AI/ML ứng dụng trong kiểm thử phần mềm, đồng thời xây dựng chương trình thực nghiệm phân loại mức độ lỗi bằng mô hình Cây quyết định (Decision Tree)."
    ]
)

# Slide 3: Chương 1: Lý do & Mục tiêu
create_content_slide(
    3,
    "Chương 1: Tổng quan",
    "Lý do chọn đề tài & Mục tiêu nghiên cứu",
    layout_type="split_image",
    bullets=[
        "**Khối lượng lỗi lớn:** Trong các dự án lớn, hàng trăm lỗi (bugs) phát sinh mỗi ngày. QA tốn nhiều thời gian để phân loại thủ công.",
        "**Ưu tiên sửa lỗi:** Việc gán mức độ nghiêm trọng (Severity) giúp xác định thứ tự ưu tiên sửa lỗi, tối ưu hóa tài nguyên và nhân sự.",
        "**Tự động hóa thông minh:** Áp dụng Machine Learning giúp tự động phân loại mức độ nghiêm trọng của lỗi nhanh chóng, khách quan.",
        "**Mục tiêu nghiên cứu:** Khảo sát các ứng dụng AI/ML trong kiểm thử, xây dựng mô hình Cây quyết định (Decision Tree) và đánh giá độ chính xác."
    ],
    image_path="bug_classification_ai.png",
    image_rect=(Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.2))
)

# Slide 4: Chương 1: Phạm vi & Phương pháp
create_content_slide(
    4,
    "Chương 1: Tổng quan",
    "Phạm vi & Phương pháp nghiên cứu",
    bullets=[
        "**Phạm vi nghiên cứu:** Tập trung tìm hiểu lý thuyết AI/ML cơ bản liên quan đến kiểm thử phần mềm. Phần thực nghiệm giới hạn trong việc xây dựng mô hình phân loại lỗi sử dụng thuật toán Cây quyết định (Decision Tree).",
        "**Bộ dữ liệu thực nghiệm:** Sử dụng tệp dữ liệu mẫu `bugs.csv` chứa các thuộc tính quan trọng: số lượng lỗi phát hiện (Bug_Count), trạng thái làm sập hệ thống (Crash) và nhãn mức độ lỗi (Severity).",
        "**Phương pháp nghiên cứu lý thuyết:** Đọc và phân tích các tài liệu học thuật về AI, các giải pháp QA hiện đại tích hợp AI/ML.",
        "**Phương pháp nghiên cứu thực nghiệm:** Lập trình Python, sử dụng thư viện xử lý dữ liệu Pandas và thư viện Scikit-learn để xây dựng, huấn luyện và kiểm thử mô hình."
    ]
)

# Slide 5: Divider Chapter 2
create_divider_slide("Chương 2", "Cơ sở lý thuyết về AI/ML trong QA")

# Slide 6: Chương 2: AI trong kiểm thử
create_content_slide(
    6,
    "Chương 2: Cơ sở lý thuyết",
    "Vai trò của Trí tuệ nhân tạo (AI) trong kiểm thử",
    bullets=[
        "**Khái niệm AI:** Là lĩnh vực khoa học máy tính nghiên cứu và phát triển các hệ thống có khả năng mô phỏng trí thông minh của con người (nhận diện mẫu, học tập, lập luận).",
        "**Tự động sinh Test Case:** AI phân tích tài liệu đặc tả, hành vi người dùng để tự động thiết kế và đề xuất kịch bản kiểm thử phù hợp.",
        "**Dự đoán lỗi (Bug Prediction):** Phân tích mã nguồn và lịch sử phát triển để chỉ ra các module/chức năng có nguy cơ phát sinh lỗi cao nhất.",
        "**Phân tích Log tự động:** Sử dụng AI để rà quét hàng triệu dòng log hệ thống, phát hiện các bất thường (anomalies) trong thời gian thực.",
        "**Kiểm thử giao diện (UI Visual Testing):** So sánh trực quan giao diện (Visual Diff) giữa các phiên bản phần mềm để tìm lỗi giao diện tự động."
    ]
)

# Slide 7: Chương 2: Machine Learning
create_content_slide(
    7,
    "Chương 2: Cơ sở lý thuyết",
    "Phân loại các phương pháp Học máy (Machine Learning)",
    bullets=[
        "**Khái niệm Machine Learning:** Nhánh con của AI giúp máy tính tự học hỏi quy luật từ dữ liệu thực tế mà không cần lập trình các quy tắc cứng.",
        "**Học có giám sát (Supervised Learning):** Dữ liệu huấn luyện có đầy đủ thuộc tính (Input) và nhãn kết quả (Label/Target). Mô hình học cách dự đoán nhãn cho dữ liệu mới. *Đây là phương pháp chính được sử dụng trong đề tài (nhãn: Severity).* ",
        "**Học không giám sát (Unsupervised Learning):** Dữ liệu không có nhãn sẵn. Mô hình tự gom nhóm dữ liệu (Clustering) có đặc tính giống nhau. Áp dụng gom nhóm lỗi tương tự, gom nhóm test case trùng lặp.",
        "**Học tăng cường (Reinforcement Learning):** Mô hình tự học thông qua hành động tương tác với môi trường và nhận phản hồi thưởng/phạt. Áp dụng tối ưu hóa thứ tự chạy bộ test case."
    ]
)

# Slide 8: Chương 2: Thuật toán Decision Tree
create_content_slide(
    8,
    "Chương 2: Cơ sở lý thuyết",
    "Nguyên lý hoạt động của thuật toán Cây quyết định",
    layout_type="split_image",
    bullets=[
        "**Nguyên lý hoạt động:** Là thuật toán học có giám sát. Xây dựng một cấu trúc hình cây bao gồm các nút quyết định (Decision Nodes) chứa điều kiện kiểm tra, các nhánh (Branches) và nút lá (Leaf Nodes) chứa nhãn kết quả.",
        "**Quy trình phân loại:** Bắt đầu từ nút gốc, kiểm tra giá trị của thuộc tính, rẽ nhánh tương ứng, tiếp tục kiểm tra nút con cho đến khi đạt được kết quả ở nút lá.",
        "**Luật quyết định 1 (Critical):** Nếu Crash = 1 -> Lỗi nghiêm trọng (Critical).",
        "**Luật quyết định 2 (Minor):** Nếu Crash = 0 và Bug_Count <= 3 -> Lỗi nhẹ (Minor).",
        "**Luật quyết định 3 (Major):** Nếu Crash = 0 và Bug_Count > 3 -> Lỗi trung bình (Major).",
        "**Đặc tính thuật toán:** Thuật toán rất trực quan, dễ lập trình, thời gian chạy cực kỳ nhanh, phù hợp dữ liệu kích thước nhỏ."
    ],
    image_path="decision_tree_flow.png",
    image_rect=(Inches(6.8), Inches(2.0), Inches(5.6), Inches(4.2))
)

# Slide 9: Divider Chapter 3
create_divider_slide("Chương 3", "Phân tích bài toán & Thiết kế mô hình")

# Slide 10: Chương 3: Phân tích bài toán
create_content_slide(
    10,
    "Chương 3: Phân tích & Thiết kế",
    "Mô tả bài toán phân loại mức độ nghiêm trọng của lỗi",
    bullets=[
        "**Mô hình bài toán:** Phân loại đa lớp (Multi-class Classification) dựa trên các thuộc tính kỹ thuật của lỗi.",
        "**Thuộc tính Bug_Count (Số lượng lỗi):** Biểu thị tần suất xuất hiện lỗi trong chức năng (Kiểu số nguyên). Tần suất cao chỉ ra mức ảnh hưởng lớn.",
        "**Thuộc tính Crash (Gây sập hệ thống):** Trạng thái lỗi có gây sập phần mềm hay không (0: Không sập, 1: Gây sập).",
        "**Nhãn Severity (Mức độ nghiêm trọng):** Gồm 3 lớp nhãn: Minor (Lỗi nhẹ), Major (Lỗi trung bình), Critical (Lỗi nghiêm trọng)."
    ]
)

# Slide 11: Chương 3: Bộ dữ liệu thực nghiệm
slide_data = prs.slides.add_slide(blank_layout)
set_slide_background(slide_data, COLOR_BG_LIGHT)
draw_accent_bar(slide_data, dark_theme=False)
add_footer(slide_data, 11)

# Slide title
txBox = slide_data.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
tf = txBox.text_frame
pTag = tf.paragraphs[0]
pTag.text = "CHƯƠNG 3: PHÂN TÍCH & THIẾT KẾ"
pTag.font.name = "Segoe UI"
pTag.font.size = Pt(11)
pTag.font.bold = True
pTag.font.color.rgb = COLOR_ACCENT_BLUE

pTitle = tf.add_paragraph()
pTitle.text = "Bộ dữ liệu thực nghiệm (bugs.csv)"
pTitle.font.name = "Georgia"
pTitle.font.size = Pt(30)
pTitle.font.bold = True
pTitle.font.color.rgb = COLOR_PRIMARY_TEXT
draw_horizontal_line(slide_data, Inches(0.8), Inches(1.55), Inches(11.733), Inches(0.025), COLOR_ACCENT_TEAL)

# Left Column (Text description)
txBoxLeft = slide_data.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
tfLeft = txBoxLeft.text_frame
tfLeft.word_wrap = True
tfLeft.margin_left = Inches(0)
tfLeft.margin_top = Inches(0)
parse_and_add_text(tfLeft, "**Dữ liệu mẫu:** Bộ dữ liệu bao gồm 16 bản ghi mẫu mô phỏng các trường hợp lỗi phần mềm thực tế.", font_size=15, space_after=14)
parse_and_add_text(tfLeft, "**Cấu trúc cột:** Ba cột rõ ràng bao gồm `Bug_Count`, `Crash` và `Severity` đóng vai trò nhãn.", font_size=15, space_after=14)
parse_and_add_text(tfLeft, "**Quy luật ngầm định:** Các lỗi có `Crash = 1` hầu hết là `Critical`. Các lỗi `Crash = 0` được phân loại dựa trên số lượng lỗi tích lũy `Bug_Count`.", font_size=15, space_after=14)

# Right Column (Cards visualizing bugs.csv contents)
draw_card(slide_data, Inches(6.8), Inches(1.8), Inches(5.5), Inches(1.4), "Mẫu Lỗi Nhẹ (Minor)", ["Bug_Count = 1, 2, 3", "Crash = 0 (Không sập)", "Severity = Minor (Lỗi nhẹ)"], "info")
draw_card(slide_data, Inches(6.8), Inches(3.4), Inches(5.5), Inches(1.4), "Mẫu Lỗi Trung Bình (Major)", ["Bug_Count = 4, 5, 7", "Crash = 0 hoặc 1 (Lỗi chức năng)", "Severity = Major (Trung bình)"], "default")
draw_card(slide_data, Inches(6.8), Inches(5.0), Inches(5.5), Inches(1.4), "Mẫu Lỗi Nghiêm Trọng (Critical)", ["Bug_Count = 5, 6, 7, 8, 9, 10", "Crash = 1 (Có sập hệ thống)", "Severity = Critical (Nghiêm trọng)"], "warn")


# Slide 12: Chương 3: Quy trình mô hình
create_content_slide(
    12,
    "Chương 3: Phân tích & Thiết kế",
    "Quy trình xây dựng mô hình Machine Learning",
    bullets=[
        "**Bước 1: Đọc dữ liệu (Load Dataset):** Đọc dữ liệu từ file `bugs.csv` bằng thư viện `Pandas`, lưu trữ dưới dạng bảng DataFrame.",
        "**Bước 2: Tách thuộc tính và nhãn (X/y Split):** Tách cột dữ liệu đầu vào `X = [['Bug_Count', 'Crash']]` và nhãn mong muốn `y = ['Severity']`.",
        "**Bước 3: Chia tập Train / Test (Data Splitting):** Chia bộ dữ liệu theo tỷ lệ 80% để huấn luyện (Training set) và 20% để kiểm tra (Testing set) bằng hàm `train_test_split` của Scikit-learn.",
        "**Bước 4: Huấn luyện (Model Training):** Khởi tạo mô hình cây quyết định `DecisionTreeClassifier()` và huấn luyện mô hình bằng lệnh `model.fit(X_train, y_train)`.",
        "**Bước 5: Dự đoán và Đánh giá (Evaluation):** Thực hiện dự đoán trên tập Test bằng `model.predict(X_test)` và đánh giá độ chính xác thông qua hàm `accuracy_score`."
    ]
)

# Slide 13: Divider Chapter 4
create_divider_slide("Chương 4", "Thực nghiệm & Đánh giá kết quả")

# Slide 14: Chương 4: Môi trường
create_content_slide(
    14,
    "Chương 4: Thực nghiệm",
    "Cấu hình môi trường thực nghiệm hệ thống",
    bullets=[
        "**Hệ điều hành:** Microsoft Windows 10/11 - Môi trường kiểm thử phổ biến.",
        "**Ngôn ngữ lập trình:** Python 3.9 - Ngôn ngữ hàng đầu cho phân tích dữ liệu và Học máy nhờ thư viện phong phú.",
        "**Công cụ phát triển (IDE):** Visual Studio Code - Trình soạn thảo mã nguồn nhẹ, mạnh mẽ, hỗ trợ debug trực quan.",
        "**Thư viện xử lý dữ liệu Pandas:** Hỗ trợ đọc tệp dữ liệu CSV (`bugs.csv`), hiển thị và xử lý cấu trúc dữ liệu bảng nhanh chóng.",
        "**Thư viện Scikit-learn:** Thư viện Machine Learning cốt lõi của Python, cung cấp thuật toán Decision Tree và các công cụ đo lường hiệu suất mô hình."
    ]
)

# Slide 15: Chương 4: Kết quả đánh giá
slide_results = prs.slides.add_slide(blank_layout)
set_slide_background(slide_results, COLOR_BG_LIGHT)
draw_accent_bar(slide_results, dark_theme=False)
add_footer(slide_results, 15)

# Slide title
txBox = slide_results.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
tf = txBox.text_frame
pTag = tf.paragraphs[0]
pTag.text = "CHƯƠNG 4: THỰC NGHIỆM"
pTag.font.name = "Segoe UI"
pTag.font.size = Pt(11)
pTag.font.bold = True
pTag.font.color.rgb = COLOR_ACCENT_BLUE

pTitle = tf.add_paragraph()
pTitle.text = "Đánh giá độ chính xác của mô hình AI"
pTitle.font.name = "Georgia"
pTitle.font.size = Pt(30)
pTitle.font.bold = True
pTitle.font.color.rgb = COLOR_PRIMARY_TEXT
draw_horizontal_line(slide_results, Inches(0.8), Inches(1.55), Inches(11.733), Inches(0.025), COLOR_ACCENT_TEAL)

# Left Column (Accuracies & Comments)
txBoxLeft = slide_results.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
tfLeft = txBoxLeft.text_frame
tfLeft.word_wrap = True
tfLeft.margin_left = Inches(0)
tfLeft.margin_top = Inches(0)
parse_and_add_text(tfLeft, "**Độ chính xác đạt được:** Mô hình đạt độ chính xác **75.0%** trên tập dữ liệu kiểm tra độc lập.", font_size=15, space_after=14)
parse_and_add_text(tfLeft, "**Đánh giá kết quả:** Với kích thước dữ liệu thực nghiệm cực kỳ nhỏ (16 mẫu), kết quả 75% cho thấy thuật toán đã tìm được quy luật phân loại cơ bản.", font_size=15, space_after=14)
parse_and_add_text(tfLeft, "**Hướng cải thiện:** Độ chính xác này có thể tăng tiệm cận 95-100% nếu bổ sung thêm các mẫu dữ liệu thực tế và đa dạng hóa thuộc tính đầu vào.", font_size=15, space_after=14)

# Right Column (Cards representing evaluation metric and code prediction)
draw_card(slide_results, Inches(6.8), Inches(1.8), Inches(5.5), Inches(2.2), "Chỉ số Đánh Giá (Metric)", [
    "Công thức tính Accuracy: Số mẫu dự đoán đúng / Tổng số mẫu thử nghiệm",
    "Kết quả kiểm tra: 3/4 mẫu dự đoán chính xác (Tương đương 75%)",
    "Mô hình không bị quá khớp (overfitting)"
], "success")

draw_card(slide_results, Inches(6.8), Inches(4.2), Inches(5.5), Inches(2.2), "Dự Đoán Lỗi Mới Thực Tế", [
    "Đầu vào thử nghiệm: Bug_Count = 8, Crash = 1",
    "Kết quả dự đoán của AI: Critical (Nghiêm trọng)",
    "Nhận xét: Kết quả dự đoán rất hợp lý vì tần suất lỗi cao và có làm sập hệ thống"
], "warn")


# Slide 16: Chương 4: Đánh giá & Ứng dụng
create_content_slide(
    16,
    "Chương 4: Thực nghiệm",
    "Đánh giá ưu nhược điểm & Khả năng ứng dụng",
    bullets=[
        "**Ưu điểm nổi bật:** Mô hình cực kỳ gọn nhẹ, tốc độ huấn luyện và phản hồi mili-giây. Thuật toán Cây quyết định giúp giải thích rõ ràng các quy tắc phân loại.",
        "**Hạn chế hiện tại:** Kích thước dữ liệu mẫu thực nghiệm còn nhỏ (16 mẫu), chưa phản ánh hết mọi tình huống lỗi phức tạp trong các dự án phần mềm lớn.",
        "**Khả năng ứng dụng thực tế:** Tích hợp trực tiếp vào các hệ thống Jira, GitHub Issues để tự động gợi ý mức độ nghiêm trọng dựa trên thông số lỗi."
    ]
)

# Slide 17: Divider Q&A
create_divider_slide("Q&A", "Kết luận & Giải đáp thắc mắc")

# Slide 18: Kết luận & Hướng phát triển
slide_conclusion = prs.slides.add_slide(blank_layout)
set_slide_background(slide_conclusion, COLOR_BG_LIGHT)
draw_accent_bar(slide_conclusion, dark_theme=False)
add_footer(slide_conclusion, 18)

# Slide title
txBox = slide_conclusion.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
tf = txBox.text_frame
pTag = tf.paragraphs[0]
pTag.text = "KẾT LUẬN"
pTag.font.name = "Segoe UI"
pTag.font.size = Pt(11)
pTag.font.bold = True
pTag.font.color.rgb = COLOR_ACCENT_BLUE

pTitle = tf.add_paragraph()
pTitle.text = "Tổng kết đồ án & Lời cảm ơn"
pTitle.font.name = "Georgia"
pTitle.font.size = Pt(30)
pTitle.font.bold = True
pTitle.font.color.rgb = COLOR_PRIMARY_TEXT
draw_horizontal_line(slide_conclusion, Inches(0.8), Inches(1.55), Inches(11.733), Inches(0.025), COLOR_ACCENT_TEAL)

# Left Column (Conclusions)
txBoxLeft = slide_conclusion.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
tfLeft = txBoxLeft.text_frame
tfLeft.word_wrap = True
tfLeft.margin_left = Inches(0)
tfLeft.margin_top = Inches(0)
parse_and_add_text(tfLeft, "**Đã hoàn thành:** Tìm hiểu cơ sở lý thuyết nền tảng về AI/ML trong QA, các kỹ thuật kiểm thử thông minh.", font_size=15, space_after=14)
parse_and_add_text(tfLeft, "**Mô phỏng thành công:** Xây dựng mô hình cây quyết định dự đoán mức độ nghiêm trọng của lỗi phần mềm đạt độ chính xác 75%.", font_size=15, space_after=14)
parse_and_add_text(tfLeft, "**Hướng mở rộng:** Bổ sung tập dữ liệu thực tế lớn từ các dự án, xây dựng giao diện tương tác trực tiếp (Web Tool) để người dùng sử dụng thuận tiện.", font_size=15, space_after=14)

# Right Column (Acknowledgement Card)
draw_card(slide_conclusion, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.6), "LỜI CẢM ƠN CHÂN THÀNH", [
    "Em xin chân thành cảm ơn giảng viên hướng dẫn TS. Nguyễn Tuấn Linh đã tận tình chỉ bảo, định hướng và giúp đỡ em trong suốt quá trình thực hiện đồ án.",
    "Cảm ơn các thầy cô trong Khoa Điện Tử, Bộ môn Công nghệ thông tin - Trường Đại học Kỹ thuật Công nghiệp Thái Nguyên đã truyền đạt kiến thức nền tảng quý báu.",
    "Do kiến thức bản thân và thời gian có hạn, báo cáo khó tránh khỏi thiếu sót. Em rất mong nhận được những ý kiến đóng góp quý báu từ thầy cô Hội đồng đánh giá."
], "info")


# Save Presentation
prs.save("KTPM_Tieuluan_Presentation_v3.pptx")
print("PowerPoint presentation updated successfully as 'KTPM_Tieuluan_Presentation_v3.pptx'!")
